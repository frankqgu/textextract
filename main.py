import os
import pptx   # pip install pptx or pop install python-ppxt   
from pptx import Presentation
import pymysql

# Database connections 
conn = pymysql.connect(
       host='127.0.0.1',
       user='root',
       passwd='root123',  #chage it 
       db='office_text', 
       port = 3306,
       charset="utf8")

sql = "insert into ppt_text (file_name ,page_num , page_content,load_time) values(%s,%s,%s,now())"
cursor = conn.cursor()
for root, dirs, files in os.walk("D:\\tmp"):
    for file in files:
        if file.endswith(".pptx"):
             file_name= os.path.join(root, file)
             #print(file_name)
             prs = Presentation(file_name)  #绝对路径
#prs = Presentation('D:\\tem\\test.pptx')  #Relative path
             print('start extract')
             for i,slide in enumerate(prs.slides):
    #if i == 1:  在这里可以指定提取ppt的具体页数
                page_str=''
                for shape in slide.shapes:
                      if shape.has_text_frame:
                         text_frame = shape.text_frame
                         page_str=page_str+text_frame.text+','
                         #print(str(i)+'页：'+text_frame.text)
                print(file_name+'第'+str(i)+'页：'+page_str)
                val = (file_name, i, page_str)  
                cursor.execute(sql, val) 
                conn.commit()  #save
cursor.close()           
conn.close()
print('存入成功！')
